"""
图表生成模块
支持多种图表类型，并根据模板风格自动调整样式
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from logger import info, error
from template_manager import get_color


# 图表类型映射
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,  # 使用带标记的折线图
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

# 预定义的饼图颜色（确保足够对比度）
PIE_COLORS = [
    (46, 134, 171),    # 蓝色
    (162, 59, 114),    # 紫色
    (233, 79, 55),     # 红色
    (39, 174, 96),     # 绿色
    (243, 156, 18),    # 橙色
    (52, 152, 219),    # 浅蓝
    (155, 89, 182),    # 紫罗兰
    (26, 188, 156),    # 青色
]


def create_chart(slide, chart_data, template, left, top, width, height):
    """
    在幻灯片上创建图表
    
    Args:
        slide: 幻灯片对象
        chart_data: 图表数据配置
        template: 模板配置
        left, top, width, height: 图表位置和尺寸
    
    Returns:
        Chart对象或None
    """
    try:
        chart_type_str = chart_data.get("type", "bar")
        chart_type = CHART_TYPE_MAP.get(chart_type_str)
        
        if not chart_type:
            error(f"不支持的图表类型: {chart_type_str}")
            return None
        
        # 准备图表数据
        data = CategoryChartData()
        data.categories = chart_data.get("categories", [])
        
        series_list = chart_data.get("series", [])
        for series in series_list:
            data.add_series(series.get("name", ""), series.get("values", []))
        
        # 创建图表
        x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
        chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, data).chart
        
        # 应用模板风格
        apply_chart_style(chart, chart_data, template)
        
        info(f"图表创建成功: {chart_type_str}")
        return chart
        
    except Exception as e:
        error(f"图表创建失败: {e}")
        import traceback
        error(f"详细错误: {traceback.format_exc()}")
        return None


def apply_chart_style(chart, chart_data, template):
    """
    根据模板风格设置图表样式
    """
    # 获取图表配置
    chart_config = template.get("chart", {})
    chart_style = chart_config.get("style", "default")
    
    # 检测是否深色背景
    bg_color = get_color(template, "background")
    is_dark_bg = sum(bg_color) < 384  # RGB平均值小于128视为深色
    
    # 获取颜色
    primary_color = get_color(template, "primary")
    accent_color = get_color(template, "accent")
    secondary_color = get_color(template, "secondary")
    
    # 文字颜色根据背景决定
    text_color = (240, 240, 240) if is_dark_bg else (60, 60, 60)
    
    # 设置图表标题
    chart_title = chart_data.get("title", "")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        title_font = chart.chart_title.text_frame.paragraphs[0].font
        title_font.size = Pt(chart_config.get("title_font_size", 14))
        title_font.bold = True
        title_font.color.rgb = RGBColor(*primary_color)
    
    # 根据图表类型应用样式
    chart_type_str = chart_data.get("type", "bar")
    
    if chart_type_str == "pie":
        apply_pie_style(chart, template, is_dark_bg)
    elif chart_type_str in ["line", "line_markers"]:
        apply_line_style(chart, template, is_dark_bg, primary_color, accent_color)
    elif chart_type_str == "area":
        apply_area_style(chart, template, is_dark_bg, primary_color, accent_color)
    else:
        # 柱状图等
        apply_bar_style(chart, template, is_dark_bg, chart_style, primary_color, accent_color, secondary_color)
    
    # 设置坐标轴样式（非饼图）
    if chart_type_str != "pie":
        style_axes(chart, text_color, is_dark_bg)
    
    # 设置图例
    show_legend = chart_data.get("show_legend", True)
    series_count = len(chart_data.get("series", []))
    
    if show_legend and series_count > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        # 设置图例文字颜色
        try:
            for legend_entry in chart.legend.legend_entries:
                legend_entry.font.color.rgb = RGBColor(*text_color)
                legend_entry.font.size = Pt(10)
        except:
            pass
    elif chart_type_str == "pie":
        # 饼图总是显示图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        try:
            for legend_entry in chart.legend.legend_entries:
                legend_entry.font.color.rgb = RGBColor(*text_color)
                legend_entry.font.size = Pt(10)
        except:
            pass
    else:
        chart.has_legend = False


def apply_pie_style(chart, template, is_dark_bg):
    """
    饼图样式 - 为每个扇形设置不同颜色
    """
    try:
        plot = chart.plots[0]
        series = plot.series[0]
        
        # 为每个数据点设置颜色
        for i, point in enumerate(series.points):
            # 使用预定义颜色，确保对比度
            if i < len(PIE_COLORS):
                color = PIE_COLORS[i]
            else:
                # 超出预定义颜色数量时，基于索引生成
                color = (
                    (i * 37) % 200 + 40,
                    (i * 73) % 200 + 40,
                    (i * 113) % 200 + 40
                )
            
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(*color)
            point.format.line.color.rgb = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(255, 255, 255)
            point.format.line.width = Pt(1.5)
        
        # 设置数据标签
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = False
        data_labels.show_value = False
        data_labels.font.size = Pt(11)
        data_labels.font.bold = True
        data_labels.font.color.rgb = RGBColor(255, 255, 255)  # 标签用白色，在彩色扇形上更清晰
            
    except Exception as e:
        error(f"饼图样式设置失败: {e}")


def apply_line_style(chart, template, is_dark_bg, primary_color, accent_color):
    """
    折线图样式
    """
    try:
        plot = chart.plots[0]
        text_color = (240, 240, 240) if is_dark_bg else (60, 60, 60)
        
        colors = generate_series_colors(primary_color, accent_color, len(plot.series))
        
        for i, series in enumerate(plot.series):
            color = colors[i]
            
            # 设置线条
            series.format.line.color.rgb = RGBColor(*color)
            series.format.line.width = Pt(2.5)
            
            # 设置数据点标记
            series.marker.style = 2  # 圆形标记
            series.marker.size = 8
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = RGBColor(*color)
            series.marker.format.line.color.rgb = RGBColor(255, 255, 255)
            series.marker.format.line.width = Pt(1)
            
            # 设置数据标签
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.position = XL_LABEL_POSITION.ABOVE
            data_labels.font.size = Pt(9)
            data_labels.font.color.rgb = RGBColor(*text_color)
                
    except Exception as e:
        error(f"折线图样式设置失败: {e}")


def apply_area_style(chart, template, is_dark_bg, primary_color, accent_color):
    """
    面积图样式
    """
    try:
        plot = chart.plots[0]
        
        for i, series in enumerate(plot.series):
            # 交替使用颜色
            color = primary_color if i % 2 == 0 else accent_color
            
            # 设置填充（半透明）
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*color)
            
            # 设置线条
            series.format.line.color.rgb = RGBColor(*color)
            series.format.line.width = Pt(1.5)
            
    except Exception as e:
        error(f"面积图样式设置失败: {e}")


def apply_bar_style(chart, template, is_dark_bg, chart_style, primary_color, accent_color, secondary_color):
    """
    柱状图样式
    """
    try:
        plot = chart.plots[0]
        text_color = (240, 240, 240) if is_dark_bg else (60, 60, 60)
        
        # 根据风格生成颜色
        if chart_style == "vibrant":
            colors = [primary_color, accent_color, secondary_color] * ((len(plot.series) + 2) // 3)
        elif chart_style == "neon":
            colors = [primary_color if i % 2 == 0 else accent_color for i in range(len(plot.series))]
        else:
            colors = generate_series_colors(primary_color, secondary_color, len(plot.series))
        
        for i, series in enumerate(plot.series):
            color = colors[i] if i < len(colors) else primary_color
            
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*color)
            
            # 添加边框使柱子更清晰
            if is_dark_bg:
                series.format.line.color.rgb = RGBColor(*color)
                series.format.line.width = Pt(0.5)
            else:
                series.format.line.color.rgb = RGBColor(255, 255, 255)
                series.format.line.width = Pt(1)
            
            # 设置数据标签
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            data_labels.font.size = Pt(10)
            data_labels.font.bold = True
            data_labels.font.color.rgb = RGBColor(*text_color)
                
    except Exception as e:
        error(f"柱状图样式设置失败: {e}")


def style_axes(chart, text_color, is_dark_bg):
    """
    设置坐标轴样式
    """
    try:
        # 设置类别轴（X轴）
        category_axis = chart.category_axis
        axis_line_color = (150, 150, 150) if is_dark_bg else (180, 180, 180)
        
        category_axis.format.line.color.rgb = RGBColor(*axis_line_color)
        category_axis.format.line.width = Pt(1)
        category_axis.tick_labels.font.color.rgb = RGBColor(*text_color)
        category_axis.tick_labels.font.size = Pt(10)
        
        # 设置数值轴（Y轴）
        value_axis = chart.value_axis
        value_axis.format.line.color.rgb = RGBColor(*axis_line_color)
        value_axis.format.line.width = Pt(1)
        value_axis.tick_labels.font.color.rgb = RGBColor(*text_color)
        value_axis.tick_labels.font.size = Pt(10)
        
        # 设置网格线
        value_axis.has_major_gridlines = True
        grid_color = (80, 80, 80) if is_dark_bg else (220, 220, 220)
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(*grid_color)
        value_axis.major_gridlines.format.line.width = Pt(0.5)
        
    except Exception as e:
        # 某些图表类型可能没有坐标轴
        pass


def generate_series_colors(base_color, secondary_color, count):
    """
    生成多个系列的颜色，确保每个颜色都足够鲜明
    """
    colors = []
    
    for i in range(count):
        if i == 0:
            colors.append(base_color)
        elif i == 1 and secondary_color:
            colors.append(secondary_color)
        else:
            # 基于预定义颜色生成
            if i - 2 < len(PIE_COLORS):
                colors.append(PIE_COLORS[i - 2])
            else:
                # 调整基础颜色亮度
                factor = 1 - ((i - 2) * 0.12)
                factor = max(0.4, factor)  # 确保不会太暗
                color = tuple(max(30, min(255, int(c * factor))) for c in base_color)
                colors.append(color)
    
    return colors


def validate_chart_data(chart_data):
    """
    验证图表数据是否有效
    """
    if not chart_data:
        return False, "图表数据为空"
    
    chart_type = chart_data.get("type")
    if chart_type not in CHART_TYPE_MAP:
        return False, f"不支持的图表类型: {chart_type}"
    
    categories = chart_data.get("categories", [])
    if not categories:
        return False, "缺少图表类别数据"
    
    series = chart_data.get("series", [])
    if not series:
        return False, "缺少图表系列数据"
    
    for s in series:
        if "values" not in s:
            return False, "系列数据缺少values字段"
        if len(s["values"]) != len(categories):
            return False, f"系列 '{s.get('name', '')}' 数据长度与类别不匹配"
    
    return True, None
