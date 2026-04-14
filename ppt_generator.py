from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from logger import info, success, error
from template_manager import get_current_template, get_color
from chart_generator import create_chart, validate_chart_data


def generate_ppt(data, output_file="output.pptx", template=None):
    """
    根据JSON数据和模板生成PPT文件
    
    Args:
        data: 包含PPT结构的JSON数据
        output_file: 输出文件路径
        template: 模板配置
    
    Returns:
        bool: 生成是否成功
    """
    try:
        if template is None:
            template = get_current_template()
        
        info(f"使用模板: {template.get('name', '默认模板')}")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 生成标题页
        create_title_slide(prs, data.get("title", ""), template)
        
        # 生成内容页
        slides = data.get("slides", [])
        info(f"开始生成 {len(slides)} 页幻灯片...")
        
        for i, slide_data in enumerate(slides, 1):
            create_content_slide(prs, slide_data, template, i)
        
        prs.save(output_file)
        success(f"PPT生成成功: {output_file}")
        return True
        
    except Exception as e:
        error(f"PPT生成失败: {e}")
        import traceback
        error(f"详细错误: {traceback.format_exc()}")
        return False


def create_title_slide(prs, title_text, template):
    """
    创建标题幻灯片
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_config = template.get("title_slide", {})
    layout_config = template.get("layout", {}).get("title_slide", {})
    styles = template.get("styles", {})
    
    # 设置背景
    bg_color = get_color(template, title_config.get("background_color", "background"))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
    
    # 获取装饰颜色
    decoration_color = get_color(template, "decoration")
    
    # 根据布局样式添加装饰
    decoration_style = layout_config.get("decoration_style", "bars")
    decoration_position = layout_config.get("decoration_position", "top_bottom")
    
    add_title_decoration(slide, decoration_style, decoration_position, decoration_color, template)
    
    # 添加标题
    title_y = layout_config.get("title_y", 2.5)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(title_y), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.alignment = PP_ALIGN.CENTER
    
    title_color = get_color(template, title_config.get("title_color", "primary"))
    title_para.font.size = Pt(title_config.get("title_font_size", 44))
    title_para.font.bold = title_config.get("title_font_bold", True)
    title_para.font.color.rgb = RGBColor(*title_color)


def add_title_decoration(slide, style, position, color, template):
    """
    为标题页添加装饰
    """
    color_rgb = RGBColor(*color)
    
    if style == "bars":
        # 顶部和底部装饰条
        if position in ["top_bottom", "top"]:
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.3))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = color_rgb
            top_bar.line.fill.background()
        
        if position in ["top_bottom", "bottom"]:
            bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = color_rgb
            bottom_bar.line.fill.background()
    
    elif style == "sidebar":
        # 左侧装饰条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color_rgb
        bar.line.fill.background()
    
    elif style == "geometric":
        # 几何装饰（角落方块）
        corner_size = Inches(0.6)
        # 左上角
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), corner_size, corner_size)
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = color_rgb
        shape1.line.fill.background()
        # 右下角
        shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.333) - corner_size, Inches(7.5) - corner_size, corner_size, corner_size)
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = color_rgb
        shape2.line.fill.background()
    
    elif style == "minimal":
        # 底部线条
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.5), Inches(5.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = color_rgb
        line.line.fill.background()
    
    elif style == "circles":
        # 散落的圆形装饰
        positions = [(1, 1), (11, 0.5), (0.5, 5), (10, 6)]
        for x, y in positions:
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.8), Inches(0.8))
            circle.fill.solid()
            circle.fill.fore_color.rgb = color_rgb
            circle.line.fill.background()
    
    elif style == "leaf":
        # 叶子形状装饰
        for i in range(3):
            leaf = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3 + i * 0.5), Inches(6.5), Inches(0.3), Inches(0.6))
            leaf.fill.solid()
            leaf.fill.fore_color.rgb = color_rgb
            leaf.line.fill.background()
    
    elif style == "glow":
        # 发光边框效果（暗黑模式）
        border_width = Inches(0.15)
        # 四边
        shapes = [
            (MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), border_width),  # 上
            (MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), border_width),  # 下
            (MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), border_width, Inches(7.5)),  # 左
            (MSO_SHAPE.RECTANGLE, Inches(13.183), Inches(0), border_width, Inches(7.5)),  # 右
        ]
        for shape_type, left, top, width, height in shapes:
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color_rgb
            shape.line.fill.background()


def create_content_slide(prs, slide_data, template, slide_num):
    """
    创建内容幻灯片
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    content_config = template.get("content_slide", {})
    layout_config = template.get("layout", {}).get("content_slide", {})
    styles = template.get("styles", {})
    
    # 设置背景
    bg_color = get_color(template, content_config.get("background_color", "background"))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
    
    # 获取颜色
    decoration_color = get_color(template, "decoration")
    
    # 添加装饰
    decoration_style = layout_config.get("decoration_style", "bars")
    decoration_position = layout_config.get("decoration_position", "left_top")
    add_content_decoration(slide, decoration_style, decoration_position, decoration_color, template)
    
    # 添加标题
    padding = layout_config.get("content_padding", 0.5)
    title_box = slide.shapes.add_textbox(Inches(padding), Inches(0.4), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", f"第{slide_num}页")
    
    title_color = get_color(template, content_config.get("title_color", "primary"))
    title_para.font.size = Pt(content_config.get("title_font_size", 32))
    title_para.font.bold = content_config.get("title_font_bold", True)
    title_para.font.color.rgb = RGBColor(*title_color)
    
    # 标题下划线装饰
    if decoration_style in ["bars", "line"]:
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(padding), Inches(1.15), Inches(3), Inches(0.04))
        underline.fill.solid()
        underline.fill.fore_color.rgb = RGBColor(*decoration_color)
        underline.line.fill.background()
    
    # 检查是否有图表数据
    chart_data = slide_data.get("chart")
    has_chart = chart_data and validate_chart_data(chart_data)[0]
    
    if has_chart:
        # 有图表时使用左右布局
        create_content_with_chart(slide, slide_data, template, padding, chart_data)
    else:
        # 无图表时使用全宽布局
        create_content_without_chart(slide, slide_data, template, padding)
    
    # 添加页码
    add_page_number(slide, slide_num, template)


def create_content_with_chart(slide, slide_data, template, padding, chart_data):
    """
    创建带图表的内容页（左文右图布局）
    """
    content_config = template.get("content_slide", {})
    
    # 左侧内容区域（约60%宽度）
    content_width = 6.5
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(padding), Inches(1.4), Inches(content_width), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()
    
    # 获取字体颜色
    bullet_color = get_color(template, content_config.get("bullet_color", "content"))
    content_color = get_color(template, content_config.get("content_color", "666666"))
    
    # 添加要点
    bullets = slide_data.get("bullets", [])
    if bullets:
        styles = template.get("styles", {})
        bullet_style = styles.get("bullet_style", "disc")
        bullet_prefix = get_bullet_prefix(bullet_style)
        
        for bullet in bullets:
            p = content_frame.add_paragraph()
            p.text = f"{bullet_prefix} {bullet}"
            p.font.size = Pt(content_config.get("bullet_font_size", 18))
            p.font.color.rgb = RGBColor(*bullet_color)
            p.space_after = Pt(8)
    
    # 添加详细内容
    content_text = slide_data.get("content", "")
    if content_text:
        if bullets:
            content_frame.add_paragraph()
        p = content_frame.add_paragraph()
        p.text = content_text
        p.font.size = Pt(content_config.get("content_font_size", 16))
        p.font.color.rgb = RGBColor(*content_color)
        p.font.italic = True
        p.space_after = Pt(12)
    
    # 右侧图表区域
    chart_left = padding + content_width + 0.3
    chart_top = 1.4
    chart_width = 5.5
    chart_height = 4.0
    
    create_chart(slide, chart_data, template, chart_left, chart_top, chart_width, chart_height)
    
    # 添加关键信息（放在底部）
    key_message = slide_data.get("key_message", "")
    if not key_message and bullets:
        key_message = bullets[0]
    
    if key_message:
        add_key_message_box(slide, key_message, content_config, template, padding, width=12.333 - padding * 2)


def create_content_without_chart(slide, slide_data, template, padding):
    """
    创建无图表的内容页（全宽布局）
    """
    content_config = template.get("content_slide", {})
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(padding), Inches(1.4), Inches(12.333 - padding * 2), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()
    
    # 获取字体颜色
    bullet_color = get_color(template, content_config.get("bullet_color", "content"))
    content_color = get_color(template, content_config.get("content_color", "666666"))
    accent_color = get_color(template, content_config.get("key_message_color", "accent"))
    
    # 添加要点
    bullets = slide_data.get("bullets", [])
    if bullets:
        styles = template.get("styles", {})
        bullet_style = styles.get("bullet_style", "disc")
        bullet_prefix = get_bullet_prefix(bullet_style)
        
        for bullet in bullets:
            p = content_frame.add_paragraph()
            p.text = f"{bullet_prefix} {bullet}"
            p.font.size = Pt(content_config.get("bullet_font_size", 18))
            p.font.color.rgb = RGBColor(*bullet_color)
            p.space_after = Pt(8)
        
        content_frame.add_paragraph()
    
    # 添加详细内容
    content_text = slide_data.get("content", "")
    if content_text:
        p = content_frame.add_paragraph()
        p.text = content_text
        p.font.size = Pt(content_config.get("content_font_size", 16))
        p.font.color.rgb = RGBColor(*content_color)
        p.font.italic = True
        p.space_after = Pt(12)
        content_frame.add_paragraph()
    
    # 添加关键信息
    key_message = slide_data.get("key_message", "")
    if not key_message and bullets:
        key_message = bullets[0]
    
    if key_message:
        add_key_message_box(slide, key_message, content_config, template, padding)


def add_content_decoration(slide, style, position, color, template):
    """
    为内容页添加装饰
    """
    color_rgb = RGBColor(*color)
    
    if style == "bars":
        if position in ["left_top", "left"]:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color_rgb
            bar.line.fill.background()
        if position in ["left_top", "top"]:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color_rgb
            bar.line.fill.background()
    
    elif style == "header_bar":
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color_rgb
        bar.line.fill.background()
    
    elif style == "dots":
        # 右侧边缘的点状装饰
        for i in range(5):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.8), Inches(1 + i * 1.2), Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color_rgb
            dot.line.fill.background()
    
    elif style == "line":
        # 简洁线条装饰（学术风格）
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.3), Inches(12.333), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = color_rgb
        line.line.fill.background()
    
    elif style == "wave":
        # 底部波浪装饰
        for i in range(10):
            wave = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(i * 1.4), Inches(6.8), Inches(1.5), Inches(0.5))
            wave.fill.solid()
            wave.fill.fore_color.rgb = color_rgb
            wave.line.fill.background()
    
    elif style == "vine":
        # 右侧藤蔓装饰
        for i in range(4):
            leaf = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5), Inches(1 + i * 1.5), Inches(0.3), Inches(0.5))
            leaf.fill.solid()
            leaf.fill.fore_color.rgb = color_rgb
            leaf.line.fill.background()
    
    elif style == "neon":
        # 霓虹灯效果（暗黑模式）
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color_rgb
        bar.line.fill.background()


def add_key_message_box(slide, message, config, template, padding, width=None):
    """
    添加关键信息框
    
    Args:
        slide: 幻灯片对象
        message: 关键信息文本
        config: 内容配置
        template: 模板配置
        padding: 内边距
        width: 可选宽度，默认自动计算
    """
    # 获取颜色
    accent_color = get_color(template, config.get("key_message_color", "accent"))
    
    # 计算背景框颜色（淡化处理）
    bg_r = min(255, accent_color[0] + 180) if accent_color[0] < 75 else max(0, accent_color[0] - 50)
    bg_g = min(255, accent_color[1] + 180) if accent_color[1] < 75 else max(0, accent_color[1] - 50)
    bg_b = min(255, accent_color[2] + 180) if accent_color[2] < 75 else max(0, accent_color[2] - 50)
    bg_color = (bg_r, bg_g, bg_b)
    
    # 检查是否使用圆角
    styles = template.get("styles", {})
    use_rounded = styles.get("use_rounded_corners", True)
    
    # 计算宽度
    box_width = width if width else (12.333 - padding * 2)
    
    # 添加背景框
    box_top = Inches(6.2)
    box_height = Inches(0.9)
    
    if use_rounded:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(padding), box_top, Inches(box_width), box_height)
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(padding), box_top, Inches(box_width), box_height)
    
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*bg_color)
    box.line.fill.background()
    
    # 添加文本
    text_box = slide.shapes.add_textbox(Inches(padding + 0.3), Inches(6.35), Inches(box_width - 0.6), Inches(0.6))
    text_frame = text_box.text_frame
    text_para = text_frame.paragraphs[0]
    text_para.text = f"💡 {message}"
    text_para.font.size = Pt(config.get("key_message_font_size", 20))
    text_para.font.bold = True
    text_para.font.color.rgb = RGBColor(*accent_color)


def get_bullet_prefix(style):
    """
    获取项目符号前缀
    """
    styles = {
        "disc": "●",
        "circle": "○",
        "square": "■",
        "arrow": "➤",
        "star": "★",
        "leaf": "🍃",
        "check": "✓",
        "number": "•",
        "none": ""
    }
    return styles.get(style, "●")


def add_page_number(slide, page_num, template):
    """
    添加页码
    """
    secondary_color = get_color(template, "secondary")
    
    page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.3))
    page_frame = page_box.text_frame
    page_para = page_frame.paragraphs[0]
    page_para.text = str(page_num)
    page_para.alignment = PP_ALIGN.RIGHT
    page_para.font.size = Pt(12)
    page_para.font.color.rgb = RGBColor(*secondary_color)


def generate_ppt_with_template(data, template_id, output_file="output.pptx"):
    """
    使用指定模板生成PPT
    """
    from template_manager import get_template
    
    template = get_template(template_id)
    if not template:
        error(f"模板 '{template_id}' 不存在")
        return False
    
    return generate_ppt(data, output_file, template)
